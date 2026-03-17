from flask import Flask, render_template, request, jsonify
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from transformers import pipeline
import os
from dotenv import load_dotenv
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
import docx
import csv
import pandas as pd
import pptx

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Summarizer pipeline initialization
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

app = Flask(__name__)

# Predefined folder paths
PREDEFINED_FOLDER_PATHS = [
    r"C:\Users\Srinivas V\Desktop\PES\SEM4\Project\RAG - LLM\Folder 1",
    r"C:\Users\Srinivas V\Desktop\PES\SEM4\Project\RAG - LLM\Folder 2",
    r"C:\Users\Srinivas V\Desktop\PES\SEM4\Project\RAG - LLM\Folder 3",
    r"C:\Users\Srinivas V\Desktop\PES\SEM4\Project\RAG - LLM\Folder 4"
]

# Function to extract text from PDFs
def extract_text_from_pdf(pdf_path):
    pdf_reader = PdfReader(pdf_path)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text

# Function to extract text from .docx (Word) files
def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

# Function to extract text from .txt files
def extract_text_from_txt(txt_path):
    with open(txt_path, 'r', encoding='utf-8') as file:
        return file.read()

# Function to extract text from .csv files
def extract_text_from_csv(csv_path):
    text = ""
    with open(csv_path, 'r', encoding='utf-8') as file:
        reader = csv.reader(file)
        for row in reader:
            text += " ".join(row) + "\n"
    return text

# Function to extract text from .xlsx (Excel) files
def extract_text_from_xlsx(xlsx_path):
    df = pd.read_excel(xlsx_path)
    return df.to_string()

# Function to extract text from .pptx (PowerPoint) files
def extract_text_from_pptx(pptx_path):
    presentation = pptx.Presentation(pptx_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to extract text from various file types in folders
def extract_text_from_folders(folder_paths):
    all_text_chunks = []
    seen_chunks = set()
    
    for folder_path in folder_paths:
        folder_path = folder_path.strip().strip('"')
        if os.path.isdir(folder_path):
            for file_name in os.listdir(folder_path):
                file_path = os.path.join(folder_path, file_name)
                text = ""
                
                # Extract text based on file extension
                if file_name.lower().endswith('.pdf'):
                    text = extract_text_from_pdf(file_path)
                elif file_name.lower().endswith('.docx'):
                    text = extract_text_from_docx(file_path)
                elif file_name.lower().endswith('.txt'):
                    text = extract_text_from_txt(file_path)
                elif file_name.lower().endswith('.csv'):
                    text = extract_text_from_csv(file_path)
                elif file_name.lower().endswith('.xlsx'):
                    text = extract_text_from_xlsx(file_path)
                elif file_name.lower().endswith('.pptx'):
                    text = extract_text_from_pptx(file_path)
                
                if text:
                    text_chunks = get_text_chunks(text)
                    for chunk in text_chunks:
                        if chunk not in seen_chunks:
                            seen_chunks.add(chunk)
                            all_text_chunks.append({"text": chunk, "metadata": {"file_name": file_name}})
    return all_text_chunks

# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)

# Vector store creation and saving for future use
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    texts = [chunk['text'] for chunk in text_chunks]
    metadatas = [chunk['metadata'] for chunk in text_chunks]
    vector_store = FAISS.from_texts(texts, embedding=embeddings, metadatas=metadatas)
    vector_store.save_local("faiss_index")

# Chain for question-answering
def get_conversational_chain():
    prompt_template = """
    You are a helpful assistant answering questions based on the context provided from several documents. Use the context to provide the best possible answer.

    Context:\n{context}\n
    Question: {question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Function to handle user questions
def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question, k=3)
    relevant_texts = [doc.page_content for doc in docs]
    
    if relevant_texts:
        combined_text = " ".join(relevant_texts)
        summary = summarizer(combined_text, max_length=150, min_length=30, do_sample=False)[0]['summary_text']
    else:
        summary = "No relevant information found."
    
    unique_file_names = list(set([doc.metadata.get('file_name', 'Unknown file') for doc in docs]))
    sources = ", ".join(unique_file_names)
    return summary.strip(), sources

# Preload all files and build the vector store manually
def preload_data():
    print("Preloading all files from predefined folders.")
    text_chunks = extract_text_from_folders(PREDEFINED_FOLDER_PATHS)
    get_vector_store(text_chunks)
    print("Preloading completed.")

@app.route('/')
def index():
    return render_template('phase2.html')

@app.route('/ask', methods=['POST'])
def ask():
    user_question = request.form.get('question')
    if user_question:
        answer, sources = user_input(user_question)
        return jsonify({"answer": answer, "sources": sources})
    return jsonify({"error": "Please provide a question."})

if __name__ == "__main__":
    preload_data()  # Preload data before starting the server
    app.run(debug=True)

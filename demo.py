from flask import Flask, render_template, request, jsonify
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from dotenv import load_dotenv
import google.generativeai as genai
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain_google_genai import ChatGoogleGenerativeAI
import docx
import csv
import pandas as pd
import pptx

# ------------------ CONFIG ------------------

load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

app = Flask(__name__)

# 👉 CHANGE THESE PATHS TO YOUR SYSTEM
PREDEFINED_FOLDER_PATHS = [
    r"C:\Users\monis\Downloads\Project-20260317T153620Z-1-001\Project\RAG - LLM\Folder 1",
    r"C:\Users\monis\Downloads\Project-20260317T153620Z-1-001\Project\RAG - LLM\Folder 2"
]
def extract_text_from_pdf(pdf_path):
    text = ""
    reader = PdfReader(pdf_path)
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_csv(path):
    text = ""
    with open(path, "r", encoding="utf-8") as f:
        reader = csv.reader(f)
        for row in reader:
            text += " ".join(row) + "\n"
    return text

def extract_text_from_xlsx(path):
    df = pd.read_excel(path)
    return df.to_string()

def extract_text_from_pptx(path):
    prs = pptx.Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
def get_text_chunks(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=10000,
        chunk_overlap=1000
    )
    return splitter.split_text(text)

def extract_text_from_folders(paths):
    all_chunks = []
    seen = set()

    for folder in paths:
        if not os.path.isdir(folder):
            continue

        for file in os.listdir(folder):
            path = os.path.join(folder, file)
            text = ""

            if file.endswith(".pdf"):
                text = extract_text_from_pdf(path)
            elif file.endswith(".docx"):
                text = extract_text_from_docx(path)
            elif file.endswith(".txt"):
                text = extract_text_from_txt(path)
            elif file.endswith(".csv"):
                text = extract_text_from_csv(path)
            elif file.endswith(".xlsx"):
                text = extract_text_from_xlsx(path)
            elif file.endswith(".pptx"):
                text = extract_text_from_pptx(path)

            if text:
                chunks = get_text_chunks(text)
                for chunk in chunks:
                    if chunk not in seen:
                        seen.add(chunk)
                        all_chunks.append({
                            "text": chunk,
                            "metadata": {"file_name": file}
                        })

    return all_chunks
def get_vector_store(chunks):
    embeddings = HuggingFaceEmbeddings(
        model_name="all-MiniLM-L6-v2"
    )

    texts = [c["text"] for c in chunks]
    metas = [c["metadata"] for c in chunks]

    db = FAISS.from_texts(texts, embedding=embeddings, metadatas=metas)
    db.save_local("faiss_index")
    
def get_chain():
    prompt = PromptTemplate(
        template="""
        Answer using the context below.

        Context:
        {context}

        Question:
        {question}

        Answer:
        """,
        input_variables=["context", "question"]
    )

    model = ChatGoogleGenerativeAI(
        model="gemini-1.5-flash",
        temperature=0.3
    )

    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

def user_input(question):
    embeddings = HuggingFaceEmbeddings(
        model_name="all-MiniLM-L6-v2"
    )

    db = FAISS.load_local(
        "faiss_index",
        embeddings,
    )

    docs = db.similarity_search(question, k=3)

    chain = get_chain()
    response = chain(
        {"input_documents": docs, "question": question},
        return_only_outputs=True
    )

    sources = list(set([d.metadata.get("file_name") for d in docs]))

    return response["output_text"], ", ".join(sources)

def preload_data():
    print("Loading files...")
    chunks = extract_text_from_folders(PREDEFINED_FOLDER_PATHS)
    get_vector_store(chunks)
    print("Done.")

@app.route("/")
def index():
    return render_template("phase2.html")

@app.route("/ask", methods=["POST"])
def ask():
    q = request.form.get("question")
    if q:
        ans, src = user_input(q)
        return jsonify({"answer": ans, "sources": src})
    return jsonify({"error": "No question provided"})

# ------------------ MAIN ------------------

if __name__ == "__main__":
    preload_data()
    app.run(debug=True)

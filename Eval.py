from flask import Flask, render_template, request, jsonify
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from transformers import pipeline
import os
from dotenv import load_dotenv
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
import time
import docx
import csv
import pandas as pd
import pptx
from rouge import Rouge
from nltk.translate.bleu_score import sentence_bleu, SmoothingFunction

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Summarizer pipeline initialization
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

app = Flask(__name__)

# Store selected folder paths
selected_folder_paths = []

# Store reference answers for evaluation
reference_answers = {}

# Function to load reference answers from a CSV file
def load_reference_answers_from_csv(csv_file_path):
    reference_answers = {}
    try:
        with open(csv_file_path, 'r', encoding='utf-8') as csvfile:
            reader = csv.reader(csvfile)
            next(reader)  # Skip header if there is one
            for row in reader:
                if len(row) == 2:  # Assuming the CSV file has two columns: question and answer
                    question, answer = row
                    reference_answers[question] = answer
    except FileNotFoundError:
        print(f"Error: File {csv_file_path} not found.")
    return reference_answers

# Load reference answers when the server starts
reference_answers = load_reference_answers_from_csv('reference_answers.csv')

# Function to extract text from PDFs in a folder
def extract_text_from_folders(folder_paths):
    all_text_chunks = []
    for folder_path in folder_paths:
        folder_path = folder_path.strip().strip('"')  
        if os.path.isdir(folder_path):  
            for file_name in os.listdir(folder_path):
                file_path = os.path.join(folder_path, file_name)
                text = ""
                
                # Extract text based on file extension
                if file_name.lower().endswith('.pdf'):
                    text = extract_text_from_pdf(file_path)
                elif file_name.lower().endswith('.docx'):
                    text = extract_text_from_docx(file_path)
                elif file_name.lower().endswith('.txt'):
                    text = extract_text_from_txt(file_path)
                elif file_name.lower().endswith('.csv'):
                    text = extract_text_from_csv(file_path)
                elif file_name.lower().endswith('.xlsx'):
                    text = extract_text_from_xlsx(file_path)
                elif file_name.lower().endswith('.pptx'):
                    text = extract_text_from_pptx(file_path)

                if text:
                    text_chunks = get_text_chunks(text)
                    all_text_chunks.extend([{"text": chunk, "metadata": {"file_name": file_name}} for chunk in text_chunks])
    return all_text_chunks

# Function to extract text from PDFs
def extract_text_from_pdf(pdf_path):
    pdf_reader = PdfReader(pdf_path)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text

# Functions for other file types (docx, txt, etc.) here...
def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_txt(txt_path):
    with open(txt_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_csv(csv_path):
    text = ""
    with open(csv_path, 'r', encoding='utf-8') as file:
        reader = csv.reader(file)
        for row in reader:
            text += " ".join(row) + "\n"
    return text

def extract_text_from_xlsx(xlsx_path):
    df = pd.read_excel(xlsx_path)
    return df.to_string()

def extract_text_from_pptx(pptx_path):
    presentation = pptx.Presentation(pptx_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)

# Vector store creation and saving for future use
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    texts = [chunk['text'] for chunk in text_chunks]
    metadatas = [chunk['metadata'] for chunk in text_chunks]
    vector_store = FAISS.from_texts(texts, embedding=embeddings, metadatas=metadatas)
    vector_store.save_local("faiss_index")

# Chain for question-answering
def get_conversational_chain():
    prompt_template = """
    You are a helpful assistant answering questions based on the context provided from several documents. Use the context to provide the best possible answer.

    Context:\n{context}\n
    Question: {question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

# Function to calculate evaluation metrics
def calculate_metrics(reference, generated):
    rouge = Rouge()
    scores = rouge.get_scores(generated, reference, avg=True)

    smoothing_function = SmoothingFunction().method1
    bleu_score = sentence_bleu([reference.split()], generated.split(), smoothing_function=smoothing_function)

    metrics = {
        "rouge1": scores["rouge-1"]["f"],
        "rouge2": scores["rouge-2"]["f"],
        "rougeL": scores["rouge-l"]["f"],
        "bleu": bleu_score
    }
    return metrics

# Function to handle user questions
def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    
    # Retrieve documents from the vector store based on the user question
    docs = new_db.similarity_search(user_question, k=5)  

    # Pass the documents to the chain for generating an answer
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)

    # Identify sources from the document metadata
    file_names = [doc.metadata.get('file_name', 'Unknown file') for doc in docs]
    unique_file_names = list(set(file_names))
    sources = ", ".join(unique_file_names)

    # For evaluation, check if the question exists in reference answers and calculate metrics
    reference = reference_answers.get(user_question, "")
    metrics = calculate_metrics(reference, response["output_text"]) if reference else {}

    return response["output_text"], sources, metrics

# Route to set reference answers (no need to set this if using CSV)
@app.route('/set_reference_answers', methods=['POST'])
def set_reference_answers():
    return jsonify({"message": "Reference answers are already loaded from the CSV file."})

@app.route('/')
def index():
    return render_template('Eval.html')

@app.route('/set_folders', methods=['POST'])
def set_folders():
    global selected_folder_paths
    folder_paths = request.json.get('folders', [])
    if folder_paths:
        selected_folder_paths = folder_paths
        # Preload data based on selected folders
        text_chunks = extract_text_from_folders(selected_folder_paths)
        get_vector_store(text_chunks)
        return jsonify({"message": "Folders successfully set and loaded.", "folders": selected_folder_paths})
    return jsonify({"error": "No folders provided."}), 400

@app.route('/ask', methods=['POST'])
def ask():
    user_question = request.json.get('question')  # Handle JSON request
    if user_question:
        answer, sources, metrics = user_input(user_question)
        return jsonify({"answer": answer, "sources": sources, "metrics": metrics})
    return jsonify({"error": "Please provide a question."})

if __name__ == "__main__": 
    app.run(debug=True)
